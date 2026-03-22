"""Office document chunking: PDF, DOCX, PPTX.

Extracts text from binary office formats and chunks it using
format-aware strategies:
- PDF: page-based splitting with short-page merging
- DOCX: heading-based sectioning (mirrors markdown chunker)
- PPTX: slide-based splitting with speaker notes

Dependencies are lazy-imported so missing libraries produce a clear
error only when the relevant format is encountered.
"""

import logging
from typing import Optional

from ..config import Config
from ..models import Chunk, ChunkType
from .common import count_tokens, make_chunk_id

logger = logging.getLogger(__name__)

# Metadata keys that represent a single position and should be converted
# to start/end ranges when chunks are merged.
_RANGE_KEYS = {
    "page_number": ("start_page", "end_page"),
    "slide_number": ("start_slide", "end_slide"),
}


def _merge_range_metadata(meta_a: dict, meta_b: dict) -> dict:
    """Merge two chunk metadata dicts, converting positional keys to ranges.

    When merging page 3 and page 4, ``page_number`` is replaced with
    ``start_page: 3, end_page: 4`` so the merged chunk accurately
    reflects the span it covers.
    """
    merged = {**meta_a, **meta_b}

    for single_key, (start_key, end_key) in _RANGE_KEYS.items():
        val_a = meta_a.get(start_key, meta_a.get(single_key))
        val_b = meta_b.get(end_key, meta_b.get(single_key))
        if val_a is not None and val_b is not None:
            merged.pop(single_key, None)
            merged[start_key] = val_a
            merged[end_key] = val_b

    return merged


def _hard_split_text(text: str, max_tokens: int) -> list[str]:
    """Split a single oversized text block into pieces under max_tokens.

    Tries sentence boundaries first ('. '), then falls back to word
    boundaries. Guarantees every returned piece is ≤ max_tokens.
    """
    if count_tokens(text) <= max_tokens:
        return [text]

    # Try splitting at sentence boundaries
    sentences = text.replace(". ", ".\n").split("\n")
    if len(sentences) > 1:
        pieces: list[str] = []
        current: list[str] = []
        current_tokens = 0
        for sent in sentences:
            st = count_tokens(sent)
            if current_tokens + st > max_tokens and current:
                pieces.append(" ".join(current))
                current = []
                current_tokens = 0
            current.append(sent)
            current_tokens += st
        if current:
            pieces.append(" ".join(current))
        # Recursively split any still-oversized pieces at word level
        result: list[str] = []
        for p in pieces:
            if count_tokens(p) > max_tokens:
                result.extend(_hard_split_by_words(p, max_tokens))
            else:
                result.append(p)
        return result

    return _hard_split_by_words(text, max_tokens)


def _hard_split_by_words(text: str, max_tokens: int) -> list[str]:
    """Last-resort splitter: break at word boundaries."""
    words = text.split()
    pieces: list[str] = []
    current: list[str] = []
    current_tokens = 0
    for word in words:
        wt = count_tokens(word)
        if current_tokens + wt > max_tokens and current:
            pieces.append(" ".join(current))
            current = []
            current_tokens = 0
        current.append(word)
        current_tokens += wt
    if current:
        pieces.append(" ".join(current))
    return pieces


def chunk_office(
    abs_path: str,
    file_path: str,
    language: str,
    config: Config,
) -> list[Chunk]:
    """Dispatch to the appropriate office format chunker.

    Args:
        abs_path: Absolute path to the binary file.
        file_path: Relative path (from project root) for chunk metadata.
        language: One of "pdf", "docx", "pptx".
        config: Loaded configuration.

    Returns:
        List of Chunk objects. Empty if extraction fails or yields no text.
    """
    try:
        if language == "pdf":
            return _chunk_pdf(abs_path, file_path, config)
        elif language == "docx":
            return _chunk_docx(abs_path, file_path, config)
        elif language == "pptx":
            return _chunk_pptx(abs_path, file_path, config)
        else:
            logger.warning("Unknown office format: %s", language)
            return []
    except Exception as exc:
        logger.warning("Failed to process office file %s: %s", file_path, exc)
        return []


def _merge_short_chunks(
    chunks: list[Chunk],
    min_tokens: int,
) -> list[Chunk]:
    """Merge consecutive chunks that are below min_tokens.

    Combines adjacent short chunks into one, updating line numbers,
    token counts, and generating a new chunk ID.
    """
    if not chunks:
        return []

    merged: list[Chunk] = []
    buffer: Optional[Chunk] = None

    for chunk in chunks:
        if buffer is None:
            buffer = chunk
            continue

        if buffer.token_count < min_tokens or chunk.token_count < min_tokens:
            # Merge into buffer
            combined_content = buffer.content + "\n\n" + chunk.content
            combined_tokens = count_tokens(combined_content)
            merged_meta = _merge_range_metadata(buffer.metadata, chunk.metadata)
            buffer = Chunk(
                id=make_chunk_id(buffer.file_path, combined_content, buffer.start_line),
                file_path=buffer.file_path,
                start_line=buffer.start_line,
                end_line=chunk.end_line,
                content=combined_content,
                chunk_type=buffer.chunk_type,
                language=buffer.language,
                symbol_name=buffer.symbol_name,
                token_count=combined_tokens,
                metadata=merged_meta,
            )
        else:
            merged.append(buffer)
            buffer = chunk

    if buffer is not None:
        merged.append(buffer)

    return merged


def _split_text_at_paragraphs(
    text: str,
    file_path: str,
    chunk_type: ChunkType,
    language: str,
    base_line: int,
    max_tokens: int,
    min_tokens: int,
    metadata: dict,
) -> list[Chunk]:
    """Split text at paragraph boundaries (double newlines) when it exceeds max_tokens.

    Returns one or more Chunk objects.
    """
    if count_tokens(text) <= max_tokens:
        tc = count_tokens(text)
        if tc < min_tokens:
            return []
        return [Chunk(
            id=make_chunk_id(file_path, text, base_line),
            file_path=file_path,
            start_line=base_line,
            end_line=base_line,
            content=text,
            chunk_type=chunk_type,
            language=language,
            token_count=tc,
            metadata=dict(metadata),
        )]

    paragraphs = text.split("\n\n")
    chunks: list[Chunk] = []
    current_parts: list[str] = []
    current_tokens = 0

    for para in paragraphs:
        para_tokens = count_tokens(para)

        # If a single paragraph exceeds max_tokens, hard-split it
        if para_tokens > max_tokens:
            # Flush current buffer first
            if current_parts:
                chunk_text = "\n\n".join(current_parts)
                tc = count_tokens(chunk_text)
                if tc >= min_tokens:
                    chunks.append(Chunk(
                        id=make_chunk_id(file_path, chunk_text, base_line),
                        file_path=file_path,
                        start_line=base_line,
                        end_line=base_line,
                        content=chunk_text,
                        chunk_type=chunk_type,
                        language=language,
                        token_count=tc,
                        metadata=dict(metadata),
                    ))
                current_parts = []
                current_tokens = 0

            for sub_piece in _hard_split_text(para, max_tokens):
                tc = count_tokens(sub_piece)
                if tc >= min_tokens:
                    chunks.append(Chunk(
                        id=make_chunk_id(file_path, sub_piece, base_line),
                        file_path=file_path,
                        start_line=base_line,
                        end_line=base_line,
                        content=sub_piece,
                        chunk_type=chunk_type,
                        language=language,
                        token_count=tc,
                        metadata=dict(metadata),
                    ))
            continue

        if current_tokens + para_tokens > max_tokens and current_parts:
            chunk_text = "\n\n".join(current_parts)
            tc = count_tokens(chunk_text)
            if tc >= min_tokens:
                chunks.append(Chunk(
                    id=make_chunk_id(file_path, chunk_text, base_line),
                    file_path=file_path,
                    start_line=base_line,
                    end_line=base_line,
                    content=chunk_text,
                    chunk_type=chunk_type,
                    language=language,
                    token_count=tc,
                    metadata=dict(metadata),
                ))
            current_parts = []
            current_tokens = 0

        current_parts.append(para)
        current_tokens += para_tokens

    if current_parts:
        chunk_text = "\n\n".join(current_parts)
        tc = count_tokens(chunk_text)
        if tc >= min_tokens:
            chunks.append(Chunk(
                id=make_chunk_id(file_path, chunk_text, base_line),
                file_path=file_path,
                start_line=base_line,
                end_line=base_line,
                content=chunk_text,
                chunk_type=chunk_type,
                language=language,
                token_count=tc,
                metadata=dict(metadata),
            ))

    return chunks


def _extract_table_text(table) -> str:
    """Extract text from a table object (python-docx or python-pptx).

    Concatenates cell contents row by row, separated by ' | '.
    Rows separated by newlines.
    """
    rows: list[str] = []
    for row in table.rows:
        cells = [cell.text.strip() for cell in row.cells]
        row_text = " | ".join(cells)
        if row_text.strip():
            rows.append(row_text)
    return "\n".join(rows)


# ---------------------------------------------------------------------------
# PDF chunker
# ---------------------------------------------------------------------------

def _chunk_pdf(abs_path: str, file_path: str, config: Config) -> list[Chunk]:
    """Extract text from PDF pages and chunk by page boundaries.

    Uses PyMuPDF (fitz) for text extraction. Short consecutive pages
    are merged; long pages are split at paragraph boundaries.
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.error(
            "PyMuPDF is required for PDF indexing. "
            "Install with: pip install PyMuPDF>=1.24.0 "
            "or run: bash setup.sh --with-office"
        )
        return []

    max_tokens = config.chunking.max_tokens
    min_tokens = config.chunking.min_tokens

    try:
        doc = fitz.open(abs_path)
    except Exception as exc:
        logger.warning("Cannot open PDF %s: %s", file_path, exc)
        return []

    total_pages = len(doc)

    # Extract document metadata
    pdf_meta = doc.metadata or {}
    doc_title = pdf_meta.get("title", "") or ""
    doc_author = pdf_meta.get("author", "") or ""
    base_metadata = {"total_pages": total_pages}
    if doc_title:
        base_metadata["pdf_title"] = doc_title
    if doc_author:
        base_metadata["pdf_author"] = doc_author

    # Extract text per page
    page_texts: list[tuple[int, str]] = []  # (1-based page number, text)
    for i in range(total_pages):
        page = doc[i]
        text = page.get_text("text").strip()
        if text:
            page_texts.append((i + 1, text))

    doc.close()

    if not page_texts:
        logger.info("PDF %s has no extractable text (scanned?)", file_path)
        return []

    # Build raw chunks per page, then merge short ones
    raw_chunks: list[Chunk] = []
    for page_num, text in page_texts:
        meta = {**base_metadata, "page_number": page_num}
        page_chunks = _split_text_at_paragraphs(
            text, file_path, ChunkType.PDF_PAGE, "pdf",
            base_line=page_num, max_tokens=max_tokens,
            min_tokens=0,  # don't filter yet, merge first
            metadata=meta,
        )
        # If split produced nothing (empty after strip), create one chunk
        if not page_chunks:
            tc = count_tokens(text)
            page_chunks = [Chunk(
                id=make_chunk_id(file_path, text, page_num),
                file_path=file_path,
                start_line=page_num,
                end_line=page_num,
                content=text,
                chunk_type=ChunkType.PDF_PAGE,
                language="pdf",
                token_count=tc,
                metadata=meta,
            )]
        raw_chunks.extend(page_chunks)

    return _merge_short_chunks(raw_chunks, min_tokens)


# ---------------------------------------------------------------------------
# DOCX chunker
# ---------------------------------------------------------------------------

def _chunk_docx(abs_path: str, file_path: str, config: Config) -> list[Chunk]:
    """Extract text from DOCX by heading-based sections.

    Mirrors the markdown chunker: headings define section boundaries,
    body paragraphs are grouped under their nearest heading.
    Tables are extracted as row-by-row text.
    """
    try:
        import docx  # python-docx
    except ImportError:
        logger.error(
            "python-docx is required for DOCX indexing. "
            "Install with: pip install python-docx>=1.1.0 "
            "or run: bash setup.sh --with-office"
        )
        return []

    max_tokens = config.chunking.max_tokens
    min_tokens = config.chunking.min_tokens

    try:
        document = docx.Document(abs_path)
    except Exception as exc:
        logger.warning("Cannot open DOCX %s: %s", file_path, exc)
        return []

    # Extract core properties for metadata
    props = document.core_properties
    doc_title = props.title or ""
    doc_author = props.author or ""
    base_metadata: dict = {}
    if doc_title:
        base_metadata["docx_title"] = doc_title
    if doc_author:
        base_metadata["docx_author"] = doc_author

    # Heading style name prefix → level
    def _heading_level(style_name: str) -> Optional[int]:
        if not style_name:
            return None
        name = style_name.lower()
        if name.startswith("heading"):
            rest = name.replace("heading", "").strip()
            if rest.isdigit():
                return int(rest)
        return None

    # Walk paragraphs + tables in document order, grouping by heading
    sections: list[dict] = []  # {heading, level, heading_path, parts}
    current_heading = ""
    current_level = 0
    current_parts: list[str] = []
    heading_stack: list[str] = []  # for heading_path
    section_index = 0

    def _flush_section() -> None:
        nonlocal current_parts, section_index
        if not current_parts:
            return
        text = "\n\n".join(current_parts)
        if text.strip():
            sections.append({
                "heading": current_heading,
                "level": current_level,
                "heading_path": list(heading_stack),
                "text": text,
                "section_index": section_index,
            })
            section_index += 1
        current_parts = []

    # Iterate body elements to preserve paragraph/table order
    for element in document.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # It's a paragraph
            from docx.text.paragraph import Paragraph
            para = Paragraph(element, document)
            style_name = para.style.name if para.style else ""
            level = _heading_level(style_name)

            if level is not None:
                # New heading — flush previous section
                _flush_section()
                current_heading = para.text.strip()
                current_level = level
                # Update heading stack
                while heading_stack and len(heading_stack) >= level:
                    heading_stack.pop()
                heading_stack.append(current_heading)
                current_parts = [current_heading]
            else:
                text = para.text.strip()
                if text:
                    current_parts.append(text)

        elif tag == "tbl":
            # It's a table
            from docx.table import Table
            table = Table(element, document)
            table_text = _extract_table_text(table)
            if table_text:
                current_parts.append(table_text)

    _flush_section()

    if not sections:
        return []

    # Convert sections to chunks
    chunks: list[Chunk] = []
    for sec in sections:
        meta = {
            **base_metadata,
            "heading_path": sec["heading_path"],
            "heading_level": sec["level"],
        }
        sec_chunks = _split_text_at_paragraphs(
            sec["text"], file_path, ChunkType.DOCX_SECTION, "docx",
            base_line=sec["section_index"] + 1,
            max_tokens=max_tokens, min_tokens=min_tokens,
            metadata=meta,
        )
        chunks.extend(sec_chunks)

    return _merge_short_chunks(chunks, min_tokens)


# ---------------------------------------------------------------------------
# PPTX chunker
# ---------------------------------------------------------------------------

def _chunk_pptx(abs_path: str, file_path: str, config: Config) -> list[Chunk]:
    """Extract text from PPTX slides and chunk per slide.

    Includes text from all shapes (text frames, tables, groups)
    and speaker notes. Slides with no extractable text are skipped.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches  # noqa: F401 — validates import
    except ImportError:
        logger.error(
            "python-pptx is required for PPTX indexing. "
            "Install with: pip install python-pptx>=1.0.0 "
            "or run: bash setup.sh --with-office"
        )
        return []

    max_tokens = config.chunking.max_tokens
    min_tokens = config.chunking.min_tokens

    try:
        prs = Presentation(abs_path)
    except Exception as exc:
        logger.warning("Cannot open PPTX %s: %s", file_path, exc)
        return []

    total_slides = len(prs.slides)
    base_metadata = {"total_slides": total_slides}

    raw_chunks: list[Chunk] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        slide_title = ""

        # Extract text from shapes
        for shape in slide.shapes:
            if shape.has_table:
                parts.append(_extract_table_text(shape.table))
            elif shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        parts.append(text)

            # Detect title placeholder safely — python-pptx can raise
            # ValueError on non-placeholder shapes when accessing
            # placeholder_format, so we guard with is_placeholder first.
            if shape.has_text_frame and getattr(shape, "is_placeholder", False):
                try:
                    ph_fmt = shape.placeholder_format
                    if ph_fmt is not None and ph_fmt.idx == 0:
                        slide_title = shape.text_frame.text.strip()
                except (AttributeError, ValueError):
                    pass

        # Extract speaker notes
        has_notes = False
        if slide.has_notes_slide:
            notes_frame = slide.notes_slide.notes_text_frame
            if notes_frame:
                notes_text = notes_frame.text.strip()
                if notes_text:
                    parts.append(f"[Speaker notes] {notes_text}")
                    has_notes = True

        slide_text = "\n\n".join(p for p in parts if p)
        if not slide_text.strip():
            continue  # Skip image-only slides

        meta = {
            **base_metadata,
            "slide_number": slide_num,
            "has_notes": has_notes,
        }
        if slide_title:
            meta["slide_title"] = slide_title

        slide_chunks = _split_text_at_paragraphs(
            slide_text, file_path, ChunkType.PPTX_SLIDE, "pptx",
            base_line=slide_num, max_tokens=max_tokens,
            min_tokens=0,  # don't filter yet, merge first
            metadata=meta,
        )
        if not slide_chunks:
            tc = count_tokens(slide_text)
            slide_chunks = [Chunk(
                id=make_chunk_id(file_path, slide_text, slide_num),
                file_path=file_path,
                start_line=slide_num,
                end_line=slide_num,
                content=slide_text,
                chunk_type=ChunkType.PPTX_SLIDE,
                language="pptx",
                token_count=tc,
                metadata=meta,
            )]
        raw_chunks.extend(slide_chunks)

    return _merge_short_chunks(raw_chunks, min_tokens)
