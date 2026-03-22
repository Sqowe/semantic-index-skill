"""Tests for office document chunking (PDF, DOCX, PPTX).

Fixtures are created programmatically — no external sample files needed.
"""

import os
import tempfile
from unittest.mock import patch

import pytest

from scripts.lib.chunkers.office import (
    _hard_split_text,
    _merge_range_metadata,
    _merge_short_chunks,
    _split_text_at_paragraphs,
    chunk_office,
)
from scripts.lib.config import load_config
from scripts.lib.models import Chunk, ChunkType


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _make_config(tmp_path, max_tokens=512, min_tokens=20):
    """Create a minimal config for testing."""
    index_dir = tmp_path / ".index"
    index_dir.mkdir(exist_ok=True)
    config_file = index_dir / "config.json"
    config_file.write_text(
        '{"schema_version":"1.0",'
        '"embedding":{"provider":"openrouter","model":"test","dimensions":768,'
        '"batch_size":50,"query_prefix":"","document_prefix":"",'
        '"max_retries":1,"retry_delay_seconds":0.1,"device":null,'
        '"trust_remote_code":false},'
        f'"chunking":{{"max_tokens":{max_tokens},"overlap_tokens":10,"min_tokens":{min_tokens}}},'
        '"indexing":{"file_extensions":[".pdf",".docx",".pptx"],'
        '"exclude_patterns":[],"max_file_size_kb":500,'
        '"max_office_file_size_kb":50000,"respect_gitignore":false},'
        '"search":{"default_top_k":10,"default_threshold":0.3,"mode":"hybrid",'
        '"hybrid_alpha":0.7,"rerank_enabled":false,'
        '"rerank_model":"test","rerank_top_n":10}}'
    )
    return load_config(str(tmp_path))


def _make_chunk(file_path="test.pdf", content="hello", start_line=1,
                end_line=1, chunk_type=ChunkType.PDF_PAGE, language="pdf",
                token_count=5, metadata=None):
    from scripts.lib.chunkers.common import make_chunk_id
    return Chunk(
        id=make_chunk_id(file_path, content, start_line),
        file_path=file_path,
        start_line=start_line,
        end_line=end_line,
        content=content,
        chunk_type=chunk_type,
        language=language,
        token_count=token_count,
        metadata=metadata or {},
    )


# ===================================================================
# Unit tests for helper functions
# ===================================================================

class TestMergeRangeMetadata:
    """Tests for _merge_range_metadata (issue #2 fix)."""

    def test_converts_page_number_to_range(self):
        a = {"page_number": 3, "total_pages": 10}
        b = {"page_number": 4, "total_pages": 10}
        result = _merge_range_metadata(a, b)
        assert result["start_page"] == 3
        assert result["end_page"] == 4
        assert "page_number" not in result

    def test_converts_slide_number_to_range(self):
        a = {"slide_number": 1, "total_slides": 5}
        b = {"slide_number": 2, "total_slides": 5}
        result = _merge_range_metadata(a, b)
        assert result["start_slide"] == 1
        assert result["end_slide"] == 2
        assert "slide_number" not in result

    def test_preserves_non_range_keys(self):
        a = {"page_number": 1, "pdf_title": "Doc"}
        b = {"page_number": 2, "pdf_author": "Author"}
        result = _merge_range_metadata(a, b)
        assert result["pdf_title"] == "Doc"
        assert result["pdf_author"] == "Author"

    def test_already_ranged_metadata_extends(self):
        a = {"start_page": 1, "end_page": 2, "total_pages": 10}
        b = {"page_number": 3, "total_pages": 10}
        result = _merge_range_metadata(a, b)
        assert result["start_page"] == 1
        assert result["end_page"] == 3


class TestHardSplitText:
    """Tests for _hard_split_text (issue #3 fix)."""

    def test_short_text_returns_single(self):
        result = _hard_split_text("short text", max_tokens=100)
        assert len(result) == 1
        assert result[0] == "short text"

    def test_splits_at_sentence_boundaries(self):
        # Each sentence ~10 tokens, total ~30 tokens, max_tokens=15 forces split
        text = (
            "First sentence is fairly long with many words. "
            "Second sentence also has enough words to matter. "
            "Third sentence completes the set of text here."
        )
        result = _hard_split_text(text, max_tokens=15)
        assert len(result) > 1
        # All pieces should be non-empty
        assert all(p.strip() for p in result)

    def test_falls_back_to_word_split(self):
        # One long "sentence" with no periods
        text = " ".join(["word"] * 200)
        result = _hard_split_text(text, max_tokens=20)
        assert len(result) > 1
        from scripts.lib.chunkers.common import count_tokens
        for piece in result:
            assert count_tokens(piece) <= 20


class TestMergeShortChunks:
    """Tests for _merge_short_chunks with range metadata."""

    def test_merges_below_min_tokens(self):
        c1 = _make_chunk(content="a", token_count=3, start_line=1,
                         end_line=1, metadata={"page_number": 1})
        c2 = _make_chunk(content="b", token_count=3, start_line=2,
                         end_line=2, metadata={"page_number": 2})
        result = _merge_short_chunks([c1, c2], min_tokens=10)
        assert len(result) == 1
        assert "a" in result[0].content
        assert "b" in result[0].content

    def test_does_not_merge_above_min_tokens(self):
        c1 = _make_chunk(content="word " * 30, token_count=30,
                         metadata={"page_number": 1})
        c2 = _make_chunk(content="word " * 30, token_count=30,
                         metadata={"page_number": 2})
        result = _merge_short_chunks([c1, c2], min_tokens=10)
        assert len(result) == 2

    def test_empty_input(self):
        assert _merge_short_chunks([], min_tokens=10) == []


class TestSplitTextAtParagraphs:
    """Tests for _split_text_at_paragraphs including oversized paragraph handling."""

    def test_single_paragraph_under_limit(self):
        result = _split_text_at_paragraphs(
            "Hello world", "test.pdf", ChunkType.PDF_PAGE, "pdf",
            base_line=1, max_tokens=100, min_tokens=1, metadata={},
        )
        assert len(result) == 1

    def test_oversized_paragraph_is_split(self):
        """A single paragraph exceeding max_tokens should be hard-split."""
        big_para = " ".join(["word"] * 300)
        result = _split_text_at_paragraphs(
            big_para, "test.pdf", ChunkType.PDF_PAGE, "pdf",
            base_line=1, max_tokens=50, min_tokens=1, metadata={},
        )
        assert len(result) > 1
        from scripts.lib.chunkers.common import count_tokens
        for chunk in result:
            assert count_tokens(chunk.content) <= 50

    def test_below_min_tokens_filtered(self):
        result = _split_text_at_paragraphs(
            "hi", "test.pdf", ChunkType.PDF_PAGE, "pdf",
            base_line=1, max_tokens=100, min_tokens=100, metadata={},
        )
        assert len(result) == 0


# ===================================================================
# PDF integration tests
# ===================================================================

class TestPDFChunker:
    """Tests for _chunk_pdf using real PyMuPDF-created fixtures."""

    def test_basic_pdf_produces_chunks(self, tmp_path):
        import fitz
        config = _make_config(tmp_path)
        pdf_path = str(tmp_path / "test.pdf")
        doc = fitz.open()
        page = doc.new_page()
        text_point = fitz.Point(72, 72)
        page.insert_text(text_point, "This is page one with enough text to index properly. " * 5)
        page2 = doc.new_page()
        page2.insert_text(text_point, "This is page two with different content for searching. " * 5)
        doc.save(pdf_path)
        doc.close()

        chunks = chunk_office(pdf_path, "test.pdf", "pdf", config)
        assert len(chunks) >= 1
        assert all(c.chunk_type == ChunkType.PDF_PAGE for c in chunks)
        assert all(c.language == "pdf" for c in chunks)

    def test_empty_pdf_returns_empty(self, tmp_path):
        import fitz
        config = _make_config(tmp_path)
        pdf_path = str(tmp_path / "empty.pdf")
        doc = fitz.open()
        doc.new_page()  # blank page, no text
        doc.save(pdf_path)
        doc.close()

        chunks = chunk_office(pdf_path, "empty.pdf", "pdf", config)
        assert chunks == []

    def test_pdf_metadata_present(self, tmp_path):
        import fitz
        config = _make_config(tmp_path)
        pdf_path = str(tmp_path / "meta.pdf")
        doc = fitz.open()
        doc.set_metadata({"title": "Test Title", "author": "Test Author"})
        page = doc.new_page()
        page.insert_text(fitz.Point(72, 72), "Content for metadata test. " * 10)
        doc.save(pdf_path)
        doc.close()

        chunks = chunk_office(pdf_path, "meta.pdf", "pdf", config)
        assert len(chunks) >= 1
        assert chunks[0].metadata.get("pdf_title") == "Test Title"
        assert chunks[0].metadata.get("pdf_author") == "Test Author"


# ===================================================================
# DOCX integration tests
# ===================================================================

class TestDOCXChunker:
    """Tests for _chunk_docx using real python-docx-created fixtures."""

    def test_basic_docx_produces_chunks(self, tmp_path):
        import docx
        config = _make_config(tmp_path)
        docx_path = str(tmp_path / "test.docx")
        doc = docx.Document()
        doc.add_heading("Introduction", level=1)
        doc.add_paragraph("This is the introduction section with enough text. " * 5)
        doc.add_heading("Details", level=2)
        doc.add_paragraph("These are the details of the document content. " * 5)
        doc.save(docx_path)

        chunks = chunk_office(docx_path, "test.docx", "docx", config)
        assert len(chunks) >= 1
        assert all(c.chunk_type == ChunkType.DOCX_SECTION for c in chunks)
        assert all(c.language == "docx" for c in chunks)

    def test_docx_heading_path_metadata(self, tmp_path):
        import docx
        config = _make_config(tmp_path)
        docx_path = str(tmp_path / "headings.docx")
        doc = docx.Document()
        doc.add_heading("Chapter One", level=1)
        doc.add_heading("Section A", level=2)
        doc.add_paragraph("Content under section A with enough words to pass min tokens. " * 3)
        doc.save(docx_path)

        chunks = chunk_office(docx_path, "headings.docx", "docx", config)
        # Find the chunk with Section A content
        section_chunks = [c for c in chunks if "Section A" in c.content]
        assert len(section_chunks) >= 1
        meta = section_chunks[0].metadata
        assert "heading_path" in meta
        assert "Section A" in meta["heading_path"]

    def test_docx_with_table(self, tmp_path):
        import docx
        config = _make_config(tmp_path)
        docx_path = str(tmp_path / "table.docx")
        doc = docx.Document()
        doc.add_heading("Data Table", level=1)
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "Alpha"
        table.cell(1, 1).text = "100"
        doc.add_paragraph("Additional context after the table for indexing. " * 3)
        doc.save(docx_path)

        chunks = chunk_office(docx_path, "table.docx", "docx", config)
        assert len(chunks) >= 1
        # Table text should appear somewhere in the chunks
        all_content = " ".join(c.content for c in chunks)
        assert "Alpha" in all_content

    def test_empty_docx_returns_empty(self, tmp_path):
        import docx
        config = _make_config(tmp_path)
        docx_path = str(tmp_path / "empty.docx")
        doc = docx.Document()
        doc.save(docx_path)

        chunks = chunk_office(docx_path, "empty.docx", "docx", config)
        assert chunks == []


# ===================================================================
# PPTX integration tests
# ===================================================================

class TestPPTXChunker:
    """Tests for _chunk_pptx using real python-pptx-created fixtures."""

    def test_basic_pptx_produces_chunks(self, tmp_path):
        from pptx import Presentation
        from pptx.util import Inches
        config = _make_config(tmp_path)
        pptx_path = str(tmp_path / "test.pptx")
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # title + content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "This is slide content with enough text for indexing. " * 5
        prs.save(pptx_path)

        chunks = chunk_office(pptx_path, "test.pptx", "pptx", config)
        assert len(chunks) >= 1
        assert all(c.chunk_type == ChunkType.PPTX_SLIDE for c in chunks)
        assert all(c.language == "pptx" for c in chunks)

    def test_pptx_title_in_metadata(self, tmp_path):
        from pptx import Presentation
        config = _make_config(tmp_path)
        pptx_path = str(tmp_path / "titled.pptx")
        prs = Presentation()
        slide_layout = prs.slide_layouts[0]  # title slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "My Presentation Title"
        slide.placeholders[1].text = "Subtitle text with enough content for indexing. " * 3
        prs.save(pptx_path)

        chunks = chunk_office(pptx_path, "titled.pptx", "pptx", config)
        assert len(chunks) >= 1
        assert chunks[0].metadata.get("slide_title") == "My Presentation Title"

    def test_pptx_speaker_notes(self, tmp_path):
        from pptx import Presentation
        config = _make_config(tmp_path)
        pptx_path = str(tmp_path / "notes.pptx")
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide With Notes"
        slide.placeholders[1].text = "Main content for the slide. " * 3
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "These are speaker notes."
        prs.save(pptx_path)

        chunks = chunk_office(pptx_path, "notes.pptx", "pptx", config)
        assert len(chunks) >= 1
        all_content = " ".join(c.content for c in chunks)
        assert "speaker notes" in all_content.lower()
        # At least one chunk should have has_notes=True
        assert any(c.metadata.get("has_notes") for c in chunks)

    def test_image_only_slide_skipped(self, tmp_path):
        from pptx import Presentation
        config = _make_config(tmp_path)
        pptx_path = str(tmp_path / "imageonly.pptx")
        prs = Presentation()
        # Blank slide layout (no text placeholders)
        slide_layout = prs.slide_layouts[6]  # blank
        prs.slides.add_slide(slide_layout)
        prs.save(pptx_path)

        chunks = chunk_office(pptx_path, "imageonly.pptx", "pptx", config)
        assert chunks == []

    def test_pptx_non_placeholder_shape_no_crash(self, tmp_path):
        """Non-placeholder shapes should not crash placeholder detection (issue #1)."""
        from pptx import Presentation
        from pptx.util import Inches, Pt
        config = _make_config(tmp_path)
        pptx_path = str(tmp_path / "freeform.pptx")
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        # Add a free-form text box (not a placeholder)
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
        tf = txBox.text_frame
        tf.text = "This is a free text box, not a placeholder. " * 5
        prs.save(pptx_path)

        # Should not raise — the old hasattr code could crash here
        chunks = chunk_office(pptx_path, "freeform.pptx", "pptx", config)
        assert len(chunks) >= 1
        assert "free text box" in chunks[0].content


# ===================================================================
# Missing dependency tests
# ===================================================================

class TestMissingDependencies:
    """Test behavior when office libraries are not installed."""

    def test_missing_fitz_returns_empty(self, tmp_path):
        config = _make_config(tmp_path)
        fake_path = str(tmp_path / "fake.pdf")
        (tmp_path / "fake.pdf").write_bytes(b"%PDF-1.4 fake")
        with patch.dict("sys.modules", {"fitz": None}):
            chunks = chunk_office(fake_path, "fake.pdf", "pdf", config)
        assert chunks == []

    def test_missing_docx_returns_empty(self, tmp_path):
        config = _make_config(tmp_path)
        fake_path = str(tmp_path / "fake.docx")
        (tmp_path / "fake.docx").write_bytes(b"PK fake")
        with patch.dict("sys.modules", {"docx": None}):
            chunks = chunk_office(fake_path, "fake.docx", "docx", config)
        assert chunks == []

    def test_missing_pptx_returns_empty(self, tmp_path):
        config = _make_config(tmp_path)
        fake_path = str(tmp_path / "fake.pptx")
        (tmp_path / "fake.pptx").write_bytes(b"PK fake")
        with patch.dict("sys.modules", {"pptx": None}):
            chunks = chunk_office(fake_path, "fake.pptx", "pptx", config)
        assert chunks == []

    def test_unknown_format_returns_empty(self, tmp_path):
        config = _make_config(tmp_path)
        chunks = chunk_office("/fake/path", "test.xyz", "xyz", config)
        assert chunks == []
